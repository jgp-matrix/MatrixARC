#!/usr/bin/env python3
"""
Extract text from DOCX/PDF files and upload to Firestore knowledge base.
For DOCX: extracts text directly (no OCR needed — much faster and more accurate).
For PDF: extracts text via PyMuPDF text extraction (no OCR for text PDFs).

Usage:
  python upload_docs_to_kb.py <source_dir> <source_name> [--recursive]

Uses Firebase CLI stored credentials for Firestore access.
"""

import os
import sys
import json
import time
import urllib.request
import urllib.parse

# Config
FIRESTORE_PROJECT = "matrix-arc"
COMPANY_ID = "XODxZ8xJc0dQXGZI7jbo"
CHUNK_CHARS = 80000  # Max chars per Firestore doc (~80KB, well under 1MB limit)

def get_firebase_token():
    """Get access token from Firebase CLI stored credentials."""
    cfg_path = os.path.join(os.environ.get("HOME") or os.environ.get("USERPROFILE", ""), ".config", "configstore", "firebase-tools.json")
    with open(cfg_path) as f:
        cfg = json.load(f)
    refresh_token = cfg.get("tokens", {}).get("refresh_token")
    if not refresh_token:
        raise Exception("No refresh token in Firebase CLI config. Run: firebase login")

    data = urllib.parse.urlencode({
        "grant_type": "refresh_token",
        "refresh_token": refresh_token,
        "client_id": "563584335869-fgrhgmd47bqnekij5i8b5pr03ho849e6.apps.googleusercontent.com",
        "client_secret": "j9iVZfS8kkCEFUPaAeJV0sAi"
    }).encode()

    req = urllib.request.Request("https://oauth2.googleapis.com/token", data=data, method="POST")
    req.add_header("Content-Type", "application/x-www-form-urlencoded")
    resp = urllib.request.urlopen(req)
    return json.loads(resp.read())["access_token"]


def upload_to_firestore(token, doc_id, payload):
    """Upload a document to Firestore knowledge base."""
    url = f"https://firestore.googleapis.com/v1/projects/{FIRESTORE_PROJECT}/databases/(default)/documents/companies/{COMPANY_ID}/knowledgeBase/{doc_id}"

    fields = {}
    for k, v in payload.items():
        if isinstance(v, str):
            fields[k] = {"stringValue": v}
        elif isinstance(v, int):
            fields[k] = {"integerValue": str(v)}

    body = json.dumps({"fields": fields}).encode()
    req = urllib.request.Request(url, data=body, method="PATCH")
    req.add_header("Content-Type", "application/json")
    req.add_header("Authorization", f"Bearer {token}")

    try:
        resp = urllib.request.urlopen(req)
        print(f"  OK Uploaded {doc_id}")
    except urllib.error.HTTPError as e:
        print(f"  ERROR uploading {doc_id}: {e.code} {e.read().decode()[:200]}")


def extract_docx_text(filepath):
    """Extract text from a DOCX file."""
    from docx import Document
    doc = Document(filepath)

    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def extract_pdf_text(filepath):
    """Extract text from a PDF file using PyMuPDF."""
    import fitz
    doc = fitz.open(filepath)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(f"--- PAGE {i+1} ---\n{text}")
    doc.close()
    return "\n\n".join(pages)


def extract_pptx_text(filepath):
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        print(f"  SKIP {filepath} (python-pptx not installed)")
        return ""

    prs = Presentation(filepath)
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_text.append(" | ".join(cells))
        if slide_text:
            parts.append(f"--- SLIDE {i+1} ---\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def main():
    if len(sys.argv) < 3:
        print("Usage: python upload_docs_to_kb.py <source_dir> <source_name> [--recursive]")
        sys.exit(1)

    source_dir = sys.argv[1]
    source_name = sys.argv[2]
    recursive = "--recursive" in sys.argv

    if not os.path.isdir(source_dir):
        print(f"ERROR: Directory not found: {source_dir}")
        sys.exit(1)

    # Collect all supported files
    files = []
    if recursive:
        for root, dirs, filenames in os.walk(source_dir):
            for fn in filenames:
                ext = fn.lower().rsplit(".", 1)[-1] if "." in fn else ""
                if ext in ("docx", "pdf", "pptx"):
                    files.append(os.path.join(root, fn))
    else:
        for fn in os.listdir(source_dir):
            ext = fn.lower().rsplit(".", 1)[-1] if "." in fn else ""
            if ext in ("docx", "pdf", "pptx"):
                files.append(os.path.join(source_dir, fn))

    files.sort()
    print(f"Found {len(files)} documents in {source_dir}")
    for f in files:
        print(f"  - {os.path.basename(f)}")

    # Extract text from all files
    all_texts = []
    for filepath in files:
        fn = os.path.basename(filepath)
        ext = fn.lower().rsplit(".", 1)[-1]
        rel_path = os.path.relpath(filepath, source_dir)

        print(f"\nExtracting: {rel_path}")
        try:
            if ext == "docx":
                text = extract_docx_text(filepath)
            elif ext == "pdf":
                text = extract_pdf_text(filepath)
            elif ext == "pptx":
                text = extract_pptx_text(filepath)
            else:
                continue

            if text.strip():
                all_texts.append(f"=== {rel_path} ===\n{text}")
                print(f"  {len(text)} chars extracted")
            else:
                print(f"  (empty - skipped)")
        except Exception as e:
            print(f"  ERROR: {e}")

    total_chars = sum(len(t) for t in all_texts)
    print(f"\nTotal extracted: {len(all_texts)} documents, {total_chars} chars")

    if not all_texts:
        print("Nothing to upload.")
        sys.exit(0)

    # Get Firebase token
    print("\nAuthenticating with Firebase...")
    token = get_firebase_token()
    print("OK Got access token")

    # Chunk and upload
    combined = "\n\n".join(all_texts)
    source_key = source_name.lower().replace(" ", "_").replace(".", "")

    # Split into chunks
    chunks = []
    current_chunk = ""
    for text_block in all_texts:
        if len(current_chunk) + len(text_block) > CHUNK_CHARS and current_chunk:
            chunks.append(current_chunk)
            current_chunk = text_block
        else:
            current_chunk += ("\n\n" if current_chunk else "") + text_block
    if current_chunk:
        chunks.append(current_chunk)

    print(f"\nUploading {len(chunks)} chunks to Firestore...")

    for i, chunk in enumerate(chunks):
        doc_id = f"{source_key}_part_{i+1:03d}"
        print(f"  Chunk {i+1}/{len(chunks)}: {len(chunk)} chars")

        upload_to_firestore(token, doc_id, {
            "title": f"{source_name} — Part {i+1}",
            "source": source_name,
            "createdAt": int(time.time() * 1000),
            "content": chunk
        })

    print(f"\nOK All done! {source_name} uploaded to Firestore knowledge base.")


if __name__ == "__main__":
    main()
